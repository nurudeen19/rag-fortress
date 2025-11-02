"""
Document loader service.
Loads documents from the knowledge base directory.
"""

from pathlib import Path
from typing import Optional, Dict, Any, List
from enum import Enum

from app.config.settings import settings
from app.core.exceptions import DocumentError


class DocumentType(str, Enum):
    """Supported document types."""
    TXT = "txt"
    PDF = "pdf"
    DOCX = "docx"
    MARKDOWN = "md"
    JSON = "json"
    CSV = "csv"
    XLSX = "xlsx"
    PPTX = "pptx"


class Document:
    """Represents a loaded document."""
    
    def __init__(
        self,
        content: str,
        doc_type: DocumentType,
        metadata: Dict[str, Any],
        source: str
    ):
        self.content = content
        self.doc_type = doc_type
        self.metadata = metadata
        self.source = source


class DocumentLoader:
    """
    Loads documents from the knowledge base directory.
    
    Single Responsibility: Load documents from the centralized knowledge base.
    All documents must be placed in the KNOWLEDGE_BASE_DIR for the loader to access them.
    """
    
    def __init__(self, knowledge_base_dir: Optional[str] = None):
        """
        Initialize document loader.
        
        Args:
            knowledge_base_dir: Override default knowledge base directory
        """
        self.knowledge_base_dir = Path(knowledge_base_dir or settings.KNOWLEDGE_BASE_DIR)
        
        # Create knowledge base directory if it doesn't exist
        self.knowledge_base_dir.mkdir(parents=True, exist_ok=True)
    
    def load(self, filename: str) -> Document:
        """
        Load a document from the knowledge base directory.
        
        Args:
            filename: Name of the file (can include subdirectories within knowledge base)
            
        Returns:
            Document: Loaded document with content and metadata
            
        Raises:
            DocumentError: If file not found or unsupported type
        """
        # Resolve path relative to knowledge base
        file_path = self.knowledge_base_dir / filename
        
        if not file_path.exists():
            raise DocumentError(
                f"File not found in knowledge base: {filename}. "
                f"Please place documents in: {self.knowledge_base_dir}"
            )
        
        # Security: Ensure file is within knowledge base (prevent directory traversal)
        try:
            file_path.resolve().relative_to(self.knowledge_base_dir.resolve())
        except ValueError:
            raise DocumentError(f"Access denied: File must be within knowledge base directory")
        
        return self._load_from_path(file_path)
    
    def load_all(self, recursive: bool = True, file_types: Optional[List[str]] = None) -> List[Document]:
        """
        Load all documents from the knowledge base directory.
        
        Args:
            recursive: Whether to search subdirectories
            file_types: List of file extensions to load (e.g., ['pdf', 'txt']). If None, loads all supported types.
            
        Returns:
            List[Document]: List of loaded documents
        """
        documents = []
        
        # Determine which file types to load
        if file_types is None:
            extensions = [f"*.{doc_type.value}" for doc_type in DocumentType]
        else:
            extensions = [f"*.{ext}" for ext in file_types]
        
        # Find all matching files
        for extension in extensions:
            if recursive:
                files = self.knowledge_base_dir.rglob(extension)
            else:
                files = self.knowledge_base_dir.glob(extension)
            
            for file_path in files:
                try:
                    document = self._load_from_path(file_path)
                    documents.append(document)
                except DocumentError as e:
                    # Log error but continue loading other documents
                    print(f"Warning: Failed to load {file_path}: {e}")
        
        return documents
    
    def _load_from_path(self, path: Path) -> Document:
    def _load_from_path(self, path: Path) -> Document:
        """
        Load document from a given path.
        
        Args:
            path: Full path to the document
            
        Returns:
            Document: Loaded document
        """
        # Determine document type from extension
        extension = path.suffix.lower().lstrip('.')
        try:
            doc_type = DocumentType(extension)
        except ValueError:
            raise DocumentError(f"Unsupported file type: {extension}")
        
        # Load content based on type
        if doc_type == DocumentType.TXT:
            content = self._load_txt(path)
        elif doc_type == DocumentType.MARKDOWN:
            content = self._load_txt(path)  # Same as txt
        elif doc_type == DocumentType.PDF:
            content = self._load_pdf(path)
        elif doc_type == DocumentType.DOCX:
            content = self._load_docx(path)
        elif doc_type == DocumentType.JSON:
            content = self._load_json(path)
        elif doc_type == DocumentType.CSV:
            content = self._load_csv(path)
        elif doc_type == DocumentType.XLSX:
            content = self._load_xlsx(path)
        elif doc_type == DocumentType.PPTX:
            content = self._load_pptx(path)
        else:
            raise DocumentError(f"No loader implemented for: {doc_type}")
        
        # Get relative path from knowledge base
        try:
            relative_path = path.relative_to(self.knowledge_base_dir)
        except ValueError:
            relative_path = path
        
        # Basic metadata
        metadata = {
            "filename": path.name,
            "relative_path": str(relative_path),
            "file_size": path.stat().st_size,
            "file_type": doc_type.value,
        }
        
        return Document(
            content=content,
            doc_type=doc_type,
            metadata=metadata,
            source=str(relative_path)
        )
    
    @staticmethod
    def _load_txt(path: Path) -> str:
        """Load plain text file."""
        try:
            return path.read_text(encoding='utf-8')
        except UnicodeDecodeError:
            # Fallback to latin-1 if utf-8 fails
            return path.read_text(encoding='latin-1')
    
    @staticmethod
    def _load_pdf(path: Path) -> str:
        """Load PDF file."""
        try:
            from pypdf2 import PdfReader
        except ImportError:
            raise DocumentError("pypdf2 not installed. Install with: pip install pypdf2")
        
        try:
            reader = PdfReader(str(path))
            text = []
            for page in reader.pages:
                text.append(page.extract_text())
            return "\n\n".join(text)
        except Exception as e:
            raise DocumentError(f"Failed to load PDF: {e}")
    
    @staticmethod
    def _load_docx(path: Path) -> str:
        """Load DOCX file."""
        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise DocumentError("python-docx not installed. Install with: pip install python-docx")
        
        try:
            doc = DocxDocument(str(path))
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            return "\n\n".join(text)
        except Exception as e:
            raise DocumentError(f"Failed to load DOCX: {e}")
    
    @staticmethod
    def _load_json(path: Path) -> str:
        """
        Load JSON file.
        Returns as formatted string for now.
        """
        import json
        try:
            data = json.loads(path.read_text())
            # Return pretty-printed JSON
            return json.dumps(data, indent=2)
        except Exception as e:
            raise DocumentError(f"Failed to load JSON: {e}")
    
    @staticmethod
    def _load_csv(path: Path) -> str:
        """
        Load CSV file.
        Returns as formatted string for now.
        """
        try:
            import csv
            rows = []
            with open(path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                for row in reader:
                    rows.append(' | '.join(row))
            return '\n'.join(rows)
        except Exception as e:
            raise DocumentError(f"Failed to load CSV: {e}")
    
    @staticmethod
    def _load_xlsx(path: Path) -> str:
        """Load Excel file."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise DocumentError("openpyxl not installed. Install with: pip install openpyxl")
        
        try:
            wb = load_workbook(str(path))
            text = []
            for sheet in wb.worksheets:
                text.append(f"Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    text.append(' | '.join(str(cell) if cell else '' for cell in row))
            return '\n'.join(text)
        except Exception as e:
            raise DocumentError(f"Failed to load XLSX: {e}")
    
    @staticmethod
    def _load_pptx(path: Path) -> str:
        """Load PowerPoint file."""
        try:
            from pptx import Presentation
        except ImportError:
            raise DocumentError("python-pptx not installed. Install with: pip install python-pptx")
        
        try:
            prs = Presentation(str(path))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n\n'.join(text)
        except Exception as e:
            raise DocumentError(f"Failed to load PPTX: {e}")
